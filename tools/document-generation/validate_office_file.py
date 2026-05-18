#!/usr/bin/env python3
"""
validate_office_file.py — Validates Microsoft Office documents (PPTX, DOCX, XLSX)
for structural integrity and content quality issues.

Usage:
  python validate_office_file.py presentation.pptx
  python validate_office_file.py report.docx --strict
  python validate_office_file.py *.xlsx --output report.json

Exit codes: 0=valid, 1=warnings, 2=errors/invalid
"""

import argparse
import json
import sys
import zipfile
from datetime import datetime
from pathlib import Path


# ── PPTX validators ──────────────────────────────────────────────────────────

def validate_pptx(path: Path, strict: bool = False) -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return {"error": "python-pptx required: pip install python-pptx"}

    issues = []
    warnings = []

    try:
        prs = Presentation(str(path))
    except Exception as e:
        return {"valid": False, "errors": [f"Cannot open: {e}"], "warnings": []}

    slide_count = len(prs.slides)
    empty_slides = 0
    text_overflow = 0
    missing_alt_text = 0

    for i, slide in enumerate(prs.slides, 1):
        has_content = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    has_content = True
                    # Check for text overflow (very long text in small box)
                    if len(text) > 500 and shape.width and shape.width < 3000000:
                        text_overflow += 1
                        warnings.append(f"Slide {i}, shape '{shape.name}': possible text overflow")

            # Check alt text on images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                if not shape.name or shape.name.startswith("Picture"):
                    missing_alt_text += 1
                    if strict:
                        issues.append(f"Slide {i}: image missing alt text (accessibility)")
                    else:
                        warnings.append(f"Slide {i}: image missing alt text")

        if not has_content:
            empty_slides += 1
            if strict:
                issues.append(f"Slide {i}: empty slide")
            else:
                warnings.append(f"Slide {i}: empty slide")

    if slide_count == 0:
        issues.append("Presentation has no slides")

    return {
        "file_type": "PPTX",
        "slides": slide_count,
        "empty_slides": empty_slides,
        "text_overflow_warnings": text_overflow,
        "missing_alt_text": missing_alt_text,
        "valid": len(issues) == 0,
        "errors": issues,
        "warnings": warnings,
    }


def validate_docx(path: Path, strict: bool = False) -> dict:
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        return {"error": "python-docx required: pip install python-docx"}

    issues = []
    warnings = []

    try:
        doc = Document(str(path))
    except Exception as e:
        return {"valid": False, "errors": [f"Cannot open: {e}"], "warnings": []}

    para_count = len(doc.paragraphs)
    empty_paras = sum(1 for p in doc.paragraphs if not p.text.strip())
    word_count = sum(len(p.text.split()) for p in doc.paragraphs)

    # Check for missing document title
    if not doc.core_properties.title:
        warnings.append("Document has no title in properties")

    # Check heading structure
    headings = [p for p in doc.paragraphs if p.style.name.startswith("Heading")]
    if para_count > 50 and not headings:
        warnings.append(f"Long document ({para_count} paragraphs) with no headings — consider adding structure")

    # Check for broken images
    broken_images = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype and rel.is_external:
            broken_images += 1
            issues.append(f"Broken external image link: {rel.target_ref}")

    # Check for very long paragraphs
    long_paras = [i + 1 for i, p in enumerate(doc.paragraphs)
                  if len(p.text) > 1500]
    if long_paras:
        if strict:
            issues.append(f"Very long paragraphs (>1500 chars) at lines: {long_paras[:5]}")
        else:
            warnings.append(f"Very long paragraphs at lines: {long_paras[:5]}")

    return {
        "file_type": "DOCX",
        "paragraphs": para_count,
        "empty_paragraphs": empty_paras,
        "word_count": word_count,
        "headings": len(headings),
        "broken_images": broken_images,
        "valid": len(issues) == 0,
        "errors": issues,
        "warnings": warnings,
    }


def validate_xlsx(path: Path, strict: bool = False) -> dict:
    try:
        import openpyxl
    except ImportError:
        return {"error": "openpyxl required: pip install openpyxl"}

    issues = []
    warnings = []

    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception as e:
        return {"valid": False, "errors": [f"Cannot open: {e}"], "warnings": []}

    sheet_info = []
    empty_sheets = 0

    for name in wb.sheetnames:
        ws = wb[name]
        row_count = ws.max_row or 0
        col_count = ws.max_column or 0

        is_empty = row_count == 0 or (row_count == 1 and col_count <= 1)
        if is_empty:
            empty_sheets += 1
            warnings.append(f"Sheet '{name}' appears empty")

        sheet_info.append({
            "name": name,
            "rows": row_count,
            "cols": col_count,
        })

        # Check for very wide sheets (may indicate data quality issues)
        if col_count > 100:
            warnings.append(f"Sheet '{name}' has {col_count} columns — verify this is intentional")

    wb.close()

    return {
        "file_type": "XLSX",
        "sheets": len(wb.sheetnames),
        "empty_sheets": empty_sheets,
        "sheet_details": sheet_info,
        "valid": len(issues) == 0,
        "errors": issues,
        "warnings": warnings,
    }


# ── ZIP integrity check ──────────────────────────────────────────────────────

def check_zip_integrity(path: Path) -> dict:
    """Office files are ZIP archives — check basic ZIP integrity first."""
    try:
        with zipfile.ZipFile(str(path), 'r') as z:
            bad = z.testzip()
            if bad:
                return {"zip_valid": False, "bad_file": bad}
        return {"zip_valid": True}
    except zipfile.BadZipFile as e:
        return {"zip_valid": False, "error": str(e)}
    except Exception as e:
        return {"zip_valid": False, "error": str(e)}


VALIDATORS = {
    ".pptx": validate_pptx,
    ".docx": validate_docx,
    ".xlsx": validate_xlsx,
}


def validate_file(path: Path, strict: bool = False) -> dict:
    ext = path.suffix.lower()
    if ext not in VALIDATORS:
        return {"file": str(path), "error": f"Unsupported type: {ext}"}

    # ZIP integrity check first
    zip_check = check_zip_integrity(path)
    if not zip_check.get("zip_valid"):
        return {
            "file": str(path),
            "valid": False,
            "errors": [f"File is corrupt or not a valid Office file: {zip_check.get('error', zip_check.get('bad_file'))}"],
            "warnings": [],
        }

    result = VALIDATORS[ext](path, strict)
    result["file"] = str(path)
    result["size_bytes"] = path.stat().st_size
    return result


def main():
    parser = argparse.ArgumentParser(description="Validate Office documents for structural integrity")
    parser.add_argument("files", nargs="+", help="Files to validate")
    parser.add_argument("--strict", action="store_true", help="Treat warnings as errors")
    parser.add_argument("--output", help="Save results to JSON file")
    args = parser.parse_args()

    results = []
    error_count = 0
    warning_count = 0

    for pattern in args.files:
        paths = list(Path(".").glob(pattern)) if "*" in pattern else [Path(pattern)]
        for p in paths:
            if not p.exists():
                print(f"SKIP: {p} not found")
                continue

            result = validate_file(p, args.strict)
            results.append(result)

            errors = result.get("errors", [])
            warnings = result.get("warnings", [])
            error_count += len(errors)
            warning_count += len(warnings)

            status = "VALID" if result.get("valid") else "INVALID"
            print(f"[{status}] {p.name} — {len(errors)} errors, {len(warnings)} warnings")
            for e in errors:
                print(f"  ERROR: {e}")
            for w in warnings[:3]:
                print(f"  WARN:  {w}")
            if len(warnings) > 3:
                print(f"  ... ({len(warnings) - 3} more warnings)")

    print(f"\nTotal: {error_count} errors, {warning_count} warnings in {len(results)} file(s)")

    if args.output:
        out = {
            "timestamp": datetime.now().isoformat(),
            "strict": args.strict,
            "error_count": error_count,
            "warning_count": warning_count,
            "files": results,
        }
        Path(args.output).write_text(json.dumps(out, indent=2, ensure_ascii=False))
        print(f"Report saved: {args.output}")

    sys.exit(2 if error_count > 0 else (1 if warning_count > 0 else 0))


if __name__ == "__main__":
    main()
