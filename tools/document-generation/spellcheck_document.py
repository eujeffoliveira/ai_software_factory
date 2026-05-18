#!/usr/bin/env python3
"""
spellcheck_document.py — Spell-check documents (PPTX, DOCX, TXT, MD, PDF).
Supports multiple backends: LanguageTool API, pyspellchecker, phunspell.

Usage:
  python spellcheck_document.py document.pptx [--lang pt-BR] [--backend languagetool]
  python spellcheck_document.py report.docx --output results.json
  python spellcheck_document.py *.md --lang en-US

Exit codes: 0=no errors, 1=spelling errors found, 2=tool error
"""

import argparse
import json
import re
import sys
from pathlib import Path
from datetime import datetime


# ── Text extractors ──────────────────────────────────────────────────────────

def extract_pptx(path: Path) -> list[tuple[str, str]]:
    """Extract (slide_ref, text) from PPTX."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx required: pip install python-pptx")
    prs = Presentation(str(path))
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            chunks.append((f"slide_{i}", " ".join(texts)))
    return chunks


def extract_docx(path: Path) -> list[tuple[str, str]]:
    """Extract (para_ref, text) from DOCX."""
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx required: pip install python-docx")
    doc = Document(str(path))
    chunks = []
    for i, para in enumerate(doc.paragraphs, 1):
        t = para.text.strip()
        if t:
            chunks.append((f"para_{i}", t))
    return chunks


def extract_txt_md(path: Path) -> list[tuple[str, str]]:
    """Extract lines from TXT/MD files."""
    lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    chunks = []
    for i, line in enumerate(lines, 1):
        # Strip markdown syntax for MD files
        clean = re.sub(r'[#*`\[\]()_~>|]', ' ', line).strip()
        if clean:
            chunks.append((f"line_{i}", clean))
    return chunks


def extract_pdf(path: Path) -> list[tuple[str, str]]:
    """Extract pages from PDF."""
    try:
        import pypdf
        reader = pypdf.PdfReader(str(path))
        chunks = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                chunks.append((f"page_{i}", text.strip()))
        return chunks
    except ImportError:
        raise RuntimeError("pypdf required: pip install pypdf")


def extract_text(path: Path) -> list[tuple[str, str]]:
    ext = path.suffix.lower()
    if ext == ".pptx":
        return extract_pptx(path)
    elif ext == ".docx":
        return extract_docx(path)
    elif ext in (".txt", ".md"):
        return extract_txt_md(path)
    elif ext == ".pdf":
        return extract_pdf(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ── Spell checkers ───────────────────────────────────────────────────────────

def check_languagetool(text: str, lang: str) -> list[dict]:
    """Check via LanguageTool REST API (local or remote)."""
    import urllib.request
    import urllib.parse
    import urllib.error

    lt_url = "http://localhost:8081/v2/check"
    data = urllib.parse.urlencode({"text": text, "language": lang}).encode()
    req = urllib.request.Request(lt_url, data=data)
    try:
        with urllib.request.urlopen(req, timeout=10) as r:
            matches = json.loads(r.read())["matches"]
    except urllib.error.URLError:
        # Fallback to public API
        lt_url = "https://api.languagetool.org/v2/check"
        with urllib.request.urlopen(
            urllib.request.Request(lt_url, data=data), timeout=15
        ) as r:
            matches = json.loads(r.read())["matches"]

    errors = []
    for m in matches:
        ctx = m["context"]
        errors.append({
            "word": ctx["text"][ctx["offset"]:ctx["offset"] + ctx["length"]],
            "message": m["message"],
            "suggestions": [r["value"] for r in m["replacements"][:3]],
            "rule": m["rule"]["id"],
        })
    return errors


def check_pyspellchecker(text: str, lang: str) -> list[dict]:
    """Check using pyspellchecker."""
    try:
        from spellchecker import SpellChecker
    except ImportError:
        raise RuntimeError("pyspellchecker required: pip install pyspellchecker")

    lang_code = lang.split("-")[0].lower()
    spell = SpellChecker(language=lang_code)
    words = re.findall(r'\b[a-zA-ZÀ-ÿ]+\b', text)
    misspelled = spell.unknown(words)
    errors = []
    for word in misspelled:
        candidates = spell.candidates(word) or set()
        errors.append({
            "word": word,
            "message": f"Possible misspelling",
            "suggestions": list(candidates)[:3],
            "rule": "SPELL",
        })
    return errors


def check_phunspell(text: str, lang: str) -> list[dict]:
    """Check using phunspell (Hunspell binding)."""
    try:
        import phunspell
    except ImportError:
        raise RuntimeError("phunspell required: pip install phunspell")

    pspell = phunspell.Phunspell(lang.replace("-", "_"))
    words = re.findall(r'\b[a-zA-ZÀ-ÿ]+\b', text)
    errors = []
    for word in words:
        if not pspell.lookup(word):
            suggestions = list(pspell.suggest(word))[:3]
            errors.append({
                "word": word,
                "message": "Misspelling (Hunspell)",
                "suggestions": suggestions,
                "rule": "SPELL",
            })
    return errors


BACKENDS = {
    "languagetool": check_languagetool,
    "pyspellchecker": check_pyspellchecker,
    "phunspell": check_phunspell,
}


# ── Main ─────────────────────────────────────────────────────────────────────

def check_file(path: Path, lang: str, backend: str) -> dict:
    checker = BACKENDS.get(backend)
    if not checker:
        raise ValueError(f"Unknown backend: {backend}. Choose from: {list(BACKENDS)}")

    chunks = extract_text(path)
    all_errors = []

    for ref, text in chunks:
        errors = checker(text, lang)
        for e in errors:
            e["location"] = ref
            all_errors.append(e)

    return {
        "file": str(path),
        "lang": lang,
        "backend": backend,
        "chunks_checked": len(chunks),
        "error_count": len(all_errors),
        "errors": all_errors,
    }


def main():
    parser = argparse.ArgumentParser(description="Spell-check PPTX, DOCX, TXT, MD, PDF documents")
    parser.add_argument("files", nargs="+", help="Files to check")
    parser.add_argument("--lang", default="pt-BR", help="Language code (default: pt-BR)")
    parser.add_argument("--backend", default="pyspellchecker",
                        choices=list(BACKENDS), help="Spell-check backend")
    parser.add_argument("--output", help="Save results to JSON file")
    parser.add_argument("--fail-on-errors", action="store_true",
                        help="Exit 1 if any spelling errors found")
    args = parser.parse_args()

    results = []
    total_errors = 0

    for pattern in args.files:
        for p in Path(".").glob(pattern) if "*" in pattern else [Path(pattern)]:
            if not p.exists():
                print(f"SKIP: {p} not found")
                continue
            print(f"Checking {p}... ", end="", flush=True)
            try:
                result = check_file(p, args.lang, args.backend)
                results.append(result)
                total_errors += result["error_count"]
                print(f"{result['error_count']} errors")
                for e in result["errors"][:5]:
                    sug = ", ".join(e["suggestions"]) if e["suggestions"] else "—"
                    print(f"  [{e['location']}] '{e['word']}': {e['message']} → {sug}")
                if result["error_count"] > 5:
                    print(f"  ... ({result['error_count'] - 5} more errors)")
            except Exception as ex:
                print(f"ERROR: {ex}")
                results.append({"file": str(p), "error": str(ex)})

    print(f"\nTotal: {total_errors} spelling errors across {len(results)} file(s)")

    if args.output:
        out = {
            "timestamp": datetime.now().isoformat(),
            "lang": args.lang,
            "backend": args.backend,
            "total_errors": total_errors,
            "files": results,
        }
        Path(args.output).write_text(json.dumps(out, indent=2, ensure_ascii=False))
        print(f"Results saved: {args.output}")

    sys.exit(1 if args.fail_on_errors and total_errors > 0 else 0)


if __name__ == "__main__":
    main()
